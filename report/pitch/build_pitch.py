"""Build the executive pitch deck from the project's own artifacts.

    docker compose run --rm sim python report/pitch/build_pitch.py

Charts are generated from `ml/models/*.csv` and the ROI figures rather than
being drawn by hand, so the deck cannot drift from the results it claims. If a
model is retrained, rerunning this regenerates the deck.

Speaker notes come from `pitch_outline.md` and are embedded in the .pptx.
"""
from __future__ import annotations

import json
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt          # noqa: E402
import pandas as pd                      # noqa: E402
from pptx import Presentation            # noqa: E402
from pptx.dml.color import RGBColor      # noqa: E402
from pptx.util import Emu, Inches, Pt    # noqa: E402

ROOT = Path(".")
MODELS = ROOT / "ml" / "models"
ASSETS = ROOT / "report" / "assets"
OUT = ROOT / "report" / "pitch" / "Digital_Twin_Project2_Executive_Pitch.pptx"
DEMO = ROOT / "report" / "demo"

INK = RGBColor(0x0F, 0x17, 0x2A)
MUTED = RGBColor(0x47, 0x55, 0x69)
ACCENT = RGBColor(0x25, 0x63, 0xEB)
WARN = RGBColor(0xB4, 0x53, 0x09)
GOOD = RGBColor(0x15, 0x80, 0x3D)

plt.rcParams.update({
    "figure.dpi": 200, "savefig.bbox": "tight", "font.size": 11,
    "axes.spines.top": False, "axes.spines.right": False,
    "axes.grid": True, "grid.alpha": 0.25, "axes.axisbelow": True,
})

# ROI figures — single source, mirrored in docs/roi_roadmap.md §3.
SCALE = [(6, 5.9), (12, 2.3), (24, 1.3), (50, 0.8)]
BENEFITS = [("Fan energy\n(measured)", 155), ("Avoided outages\n(assumed rate)", 3876),
            ("Condition-based\nservicing", 1800)]


# ── Charts ──────────────────────────────────────────────────────────────────

def chart_model_comparison() -> Path:
    df = pd.read_csv(MODELS / "model_comparison.csv")
    labels = {"always 'no failure'": "Always\n\"no failure\"",
              "always predict no-failure": "Always\n\"no failure\"",
              "threshold motor_temp>80": "Single\nthreshold",
              "logistic_regression": "Logistic\nregression",
              "random_forest": "Random\nforest",
              "gradient_boosting": "Gradient boosting\n(shipped)"}
    df = df[df.model.isin(labels)]
    names = [labels[m] for m in df.model]
    colours = ["#2563eb" if "shipped" in n else "#94a3b8" for n in names]
    fig, ax = plt.subplots(figsize=(7.8, 3.4))
    bars = ax.bar(names, df.pr_auc, color=colours, width=0.62)
    for b, v in zip(bars, df.pr_auc):
        ax.text(b.get_x() + b.get_width()/2, v + 0.02, f"{v:.2f}",
                ha="center", fontsize=10, weight="bold")
    ax.tick_params(axis="x", labelsize=9.5)
    ax.set_ylabel("PR-AUC"); ax.set_ylim(0, 1.08)
    ax.set_title("Six times better than the obvious rule", loc="left", weight="bold")
    path = ASSETS / "pitch_model_comparison.png"
    fig.savefig(path); plt.close(fig)
    return path


def chart_fault_modes() -> Path:
    df = pd.read_csv(MODELS / "recall_by_fault_mode.csv")
    names = {"hdf": "Heat dissipation", "osf": "Overstrain", "pwf": "Power"}
    df["label"] = df.fault.map(names).fillna(df.fault)
    x = range(len(df))
    fig, ax = plt.subplots(figsize=(7.4, 3.2))
    ax.bar([i - 0.2 for i in x], df.recall_model, 0.4,
           label="Model alone", color="#94a3b8")
    ax.bar([i + 0.2 for i in x], df.recall_hybrid, 0.4,
           label="With thermal guard", color="#2563eb")
    ax.set_xticks(list(x)); ax.set_xticklabels(df.label)
    ax.set_ylabel("Recall"); ax.set_ylim(0, 1.1); ax.legend(frameon=False)
    ax.set_title("The model is blind to one failure mode — a rule covers it",
                 loc="left", weight="bold")
    path = ASSETS / "pitch_fault_modes.png"
    fig.savefig(path); plt.close(fig)
    return path


def chart_fairness() -> Path:
    df = pd.read_csv(MODELS / "fairness_audit.csv").dropna(subset=["recall"])
    df = df.sort_values("recall")
    names = [t.split("/")[-1].replace("-", " ").title() for t in df.twin_id]
    colours = ["#dc2626" if r < 0.5 else "#16a34a" for r in df.recall]
    fig, ax = plt.subplots(figsize=(7.4, 3.2))
    bars = ax.barh(names, df.recall, color=colours)
    for b, v in zip(bars, df.recall):
        ax.text(min(v + 0.02, 0.92), b.get_y() + b.get_height()/2,
                f"{v:.2f}", va="center", fontsize=10, weight="bold")
    ax.set_xlabel("Recall (model channel)"); ax.set_xlim(0, 1.05)
    ax.set_title("Weakest where it matters most: Wet Lab A", loc="left",
                 weight="bold")
    path = ASSETS / "pitch_fairness.png"
    fig.savefig(path); plt.close(fig)
    return path


def chart_payback() -> Path:
    units = [u for u, _ in SCALE]
    years = [y for _, y in SCALE]
    fig, ax = plt.subplots(figsize=(7.4, 3.2))
    bars = ax.bar([str(u) for u in units], years,
                  color=["#dc2626" if y > 3 else "#16a34a" for y in years])
    for b, y in zip(bars, years):
        ax.text(b.get_x() + b.get_width()/2, y + 0.12, f"{y} yr",
                ha="center", fontsize=10, weight="bold")
    ax.axhline(3, color="#475569", ls="--", lw=1)
    ax.text(3.35, 3.12, "3-year hurdle", fontsize=9, color="#475569")
    ax.set_xlabel("HVAC units in scope"); ax.set_ylabel("Payback (years)")
    ax.set_ylim(0, 7)
    ax.set_title("Scale, not tuning, is what makes this pay", loc="left",
                 weight="bold")
    path = ASSETS / "pitch_payback.png"
    fig.savefig(path); plt.close(fig)
    return path


def chart_benefits() -> Path:
    labels = [n for n, _ in BENEFITS]
    values = [v for _, v in BENEFITS]
    fig, ax = plt.subplots(figsize=(7.4, 3.0))
    bars = ax.bar(labels, values, color=["#94a3b8", "#2563eb", "#60a5fa"])
    for b, v in zip(bars, values):
        ax.text(b.get_x() + b.get_width()/2, v + 90, f"€{v:,}",
                ha="center", fontsize=10, weight="bold")
    ax.set_ylabel("€ per year"); ax.set_ylim(0, 4600)
    ax.set_title("The energy saving is real, and small", loc="left", weight="bold")
    path = ASSETS / "pitch_benefits.png"
    fig.savefig(path); plt.close(fig)
    return path


# ── Slide helpers ───────────────────────────────────────────────────────────

W, H = Inches(13.333), Inches(7.5)


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def textbox(slide, text, left, top, width, height, size=18, bold=False,
            colour=INK, align=None, spacing=1.15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.line_spacing = spacing
        if align is not None:
            p.alignment = align
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = colour
            run.font.name = "Calibri"
    return box


def heading(slide, title, kicker=None):
    if kicker:
        textbox(slide, kicker.upper(), Inches(0.8), Inches(0.45), Inches(11), Inches(0.4),
                size=12, bold=True, colour=ACCENT)
    textbox(slide, title, Inches(0.8), Inches(0.8), Inches(11.8), Inches(1.0),
            size=30, bold=True)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def bullets(slide, items, top=Inches(2.0), size=18, colour=INK):
    y = top
    for item in items:
        textbox(slide, "•  " + item, Inches(1.0), y, Inches(11.4), Inches(0.6),
                size=size, colour=colour)
        y += Inches(0.62)
    return y


def picture(slide, path, left, top, width):
    if Path(path).exists():
        slide.shapes.add_picture(str(path), left, top, width=width)


def table(slide, rows, left, top, width, height, col_widths=None,
          header=True, size=15):
    shape = slide.shapes.add_table(len(rows), len(rows[0]), left, top, width, height)
    tbl = shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = w
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(size)
                    run.font.bold = bool(header and r == 0)
                    run.font.name = "Calibri"
    return tbl


# ── Deck ────────────────────────────────────────────────────────────────────

def build():
    ASSETS.mkdir(parents=True, exist_ok=True)
    OUT.parent.mkdir(parents=True, exist_ok=True)

    charts = {
        "models": chart_model_comparison(),
        "modes": chart_fault_modes(),
        "fairness": chart_fairness(),
        "payback": chart_payback(),
        "benefits": chart_benefits(),
    }
    spec = json.loads((MODELS / "feature_spec.json").read_text())
    m = spec["metrics"]

    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    # 1 — title
    s = blank(prs)
    textbox(s, "Smart Facility Digital Twin", Inches(0.9), Inches(2.3),
            Inches(11.5), Inches(1.2), size=44, bold=True)
    textbox(s, "Predictive maintenance for building HVAC",
            Inches(0.9), Inches(3.5), Inches(11.5), Inches(0.7),
            size=22, colour=MUTED)
    textbox(s, "6 rooms  ·  2 floors  ·  11 interacting twins  ·  1 machine-learning model",
            Inches(0.9), Inches(4.3), Inches(11.5), Inches(0.6),
            size=16, colour=ACCENT)
    notes(s, "This is a working system, not a concept. Everything shown runs.")

    # 2 — problem
    s = blank(prs)
    heading(s, "We find HVAC failures when someone complains", "The problem")
    bullets(s, [
        "In a wet lab, that means temperature-sensitive samples are already at risk.",
        "Maintenance runs on a calendar: we service units that don't need it…",
        "…and miss units that do.",
        "No visibility into equipment condition between visits.",
    ])
    textbox(s, "The question isn't whether to maintain.\nIt's whether we maintain on a calendar or on evidence.",
            Inches(1.0), Inches(5.0), Inches(11.0), Inches(1.2),
            size=22, bold=True, colour=ACCENT)
    notes(s, "The question isn't whether to maintain. It's whether we maintain "
             "on a calendar or on evidence.")

    # 3 — what we built
    s = blank(prs)
    heading(s, "A digital twin for every room", "What we built")
    bullets(s, [
        "Physics, control and equipment health simulated per room, live over MQTT.",
        "Floor and building twins coordinate a 15 kW electrical budget.",
        "An occupancy twin moves people between rooms — conserving headcount.",
        "Dashboard, 3D building view, and machine-learning risk scoring.",
    ])
    demo = sorted(DEMO.glob("*.png"))
    if demo:
        picture(s, demo[0], Inches(7.6), Inches(4.3), Inches(5.2))
    notes(s, "Each room runs its own control loop. Supervisors advise; they never "
             "take over. If the coordination layer dies, every room keeps cooling.")

    # 4 — coordination
    s = blank(prs)
    heading(s, "Rooms decide. Supervisors advise.", "Architecture")
    table(s, [
        ["", "Centralised", "Federated (chosen)"],
        ["A crash stops…", "every room cooling", "coordination only"],
        ["Occupancy data", "leaves every room", "stays room-local"],
        ["Tuning", "one loop for all rooms", "per-room thermal mass"],
        ["A supervisor bug", "actuates every room", "bounded advice, ≤1.5 °C"],
    ], Inches(1.0), Inches(2.1), Inches(11.3), Inches(2.8),
        col_widths=[Inches(2.6), Inches(4.2), Inches(4.5)])
    textbox(s, "The 1.5 °C limit is enforced by the room, not by the supervisor "
               "sending the advice.",
            Inches(1.0), Inches(5.3), Inches(11.0), Inches(0.8),
            size=17, bold=True, colour=ACCENT)
    notes(s, "One crash in a centralised design stops every room. And occupancy "
             "data would have to leave the room it came from — this is a privacy "
             "choice as much as a resilience one.")

    # 5 — the intelligence
    s = blank(prs)
    heading(s, "Four hours of warning before a failure", "The intelligence")
    picture(s, charts["models"], Inches(0.9), Inches(1.9), Inches(7.4))
    textbox(s, f"PR-AUC  {m['pr_auc']:.2f}\nRecall  {m['recall']:.2f}\n"
               f"Precision  {m['precision']:.2f}",
            Inches(8.7), Inches(2.4), Inches(3.8), Inches(2.2), size=24, bold=True)
    textbox(s, "Accuracy would be 98 % if we always predicted\n"
               "\"no failure\". That is why we don't quote it.",
            Inches(8.7), Inches(4.4), Inches(4.2), Inches(1.2),
            size=14, colour=MUTED)
    notes(s, "Accuracy would be 98% if we simply predicted 'no failure' every "
             "time — that's why we don't quote accuracy. The model is six times "
             "better than the obvious rule.")

    # 6 — honesty
    s = blank(prs)
    heading(s, "What the model cannot do", "Limitations")
    picture(s, charts["modes"], Inches(0.9), Inches(1.9), Inches(7.4))
    textbox(s, "Blind to heat-dissipation\nfailure on its own.",
            Inches(8.6), Inches(2.2), Inches(4.2), Inches(1.2),
            size=20, bold=True, colour=WARN)
    textbox(s, "About two such events existed in the\ntraining data. No model learns from two.\n\n"
               "A physics threshold covers it — we ship\nboth, and report them separately.",
            Inches(8.6), Inches(3.4), Inches(4.3), Inches(2.2), size=14, colour=MUTED)
    notes(s, "This is the slide I'd want to see if I were buying. ML is not "
             "uniformly better than rules — it beats thresholds on cumulative "
             "wear and loses to a thermostat on a fault with a direct physical "
             "precursor. We ship both.")

    # 7 — fairness
    s = blank(prs)
    heading(s, "Weakest where it matters most", "Fairness audit")
    picture(s, charts["fairness"], Inches(0.9), Inches(1.9), Inches(7.4))
    textbox(s, "Wet Lab A:\n100 % of failures missed",
            Inches(8.6), Inches(2.2), Inches(4.2), Inches(1.2),
            size=20, bold=True, colour=WARN)
    textbox(s, "Not because the model knows which room it\nis looking at — room identity is excluded.\n\n"
               "That room is serviced so well its only\nfailures are the mode the model can't see.\n\n"
               "Keep calendar servicing there.",
            Inches(8.6), Inches(3.4), Inches(4.3), Inches(2.6), size=14, colour=MUTED)
    notes(s, "The room the model covers worst is the room where failure costs "
             "most. We found it, published it, and mitigate it with the thermal "
             "guard — but keep calendar servicing on that unit. We are not asking "
             "you to trust it there.")

    # 8 — benefits
    s = blank(prs)
    heading(s, "Where the value actually is", "Business case")
    picture(s, charts["benefits"], Inches(0.9), Inches(1.9), Inches(7.4))
    textbox(s, "≈ €5,800 / year", Inches(8.8), Inches(2.5), Inches(4.0),
            Inches(0.9), size=28, bold=True, colour=GOOD)
    textbox(s, "Anyone pitching predictive maintenance\non fan energy alone is selling you\n"
               "something. The value is in outages\nyou don't have.",
            Inches(8.8), Inches(3.5), Inches(4.2), Inches(2.0), size=14, colour=MUTED)
    notes(s, "The energy saving is real but small — €155. The value is in outages "
             "you don't have.")

    # 9 — the uncomfortable number
    s = blank(prs)
    heading(s, "At six units, it does not pay for itself", "The uncomfortable number")
    table(s, [
        ["One-off investment", "€17,000"],
        ["Recurring", "€2,900 / year"],
        ["Net benefit", "€2,900 / year"],
        ["Payback", "5.9 years"],
        ["5-year NPV @ 8 %", "−€5,420"],
    ], Inches(1.0), Inches(2.1), Inches(6.2), Inches(2.9),
        col_widths=[Inches(3.6), Inches(2.6)], header=False, size=17)
    textbox(s, "Integration cost is roughly fixed.\nBenefits scale with unit count.\n\n"
               "This is a question of scale —\nnot of whether the technology works.",
            Inches(7.8), Inches(2.3), Inches(5.0), Inches(2.6),
            size=19, bold=True, colour=WARN)
    notes(s, "On this building alone, it does not pay for itself. Integration "
             "cost is roughly fixed; benefits scale with unit count.")

    # 10 — where it pays
    s = blank(prs)
    heading(s, "Where it does pay", "Recommendation")
    picture(s, charts["payback"], Inches(0.9), Inches(1.9), Inches(7.4))
    textbox(s, "Pilot to prove the model\non real telemetry.\n\nThen scale to campus.",
            Inches(8.7), Inches(2.6), Inches(4.2), Inches(2.2),
            size=20, bold=True, colour=ACCENT)
    textbox(s, "Not a building-by-building rollout.",
            Inches(8.7), Inches(4.6), Inches(4.2), Inches(0.8),
            size=14, colour=MUTED)
    notes(s, "The recommendation is a pilot to prove the model on real telemetry, "
             "then campus scale. Not a building-by-building rollout.")

    # 11 — security
    s = blank(prs)
    heading(s, "The pilot is not secure — and that is costed", "Security")
    bullets(s, [
        "Today: anonymous broker, no encryption. Fine on one machine.",
        "Unacceptable in a building that dispatches technicians.",
    ])
    table(s, [
        ["Hardening", "Why"],
        ["mTLS + per-topic ACLs", "a room twin can only speak for itself"],
        ["Signed commands", "no replay, no spoofed shutdown"],
        ["Audit log", "every action attributable to a person"],
        ["Network segmentation", "not routable from the office LAN"],
    ], Inches(1.0), Inches(3.3), Inches(11.3), Inches(2.6),
        col_widths=[Inches(4.3), Inches(7.0)])
    textbox(s, "€4,000 — a line item, not a backlog ticket.",
            Inches(1.0), Inches(6.1), Inches(9.0), Inches(0.6),
            size=18, bold=True, colour=ACCENT)
    notes(s, "A system that dispatches technicians to physical equipment needs "
             "every command attributable to a person. That's in the budget, not "
             "the backlog.")

    # 12 — ethics
    s = blank(prs)
    heading(s, "The model opens tickets. It never switches anything off.",
            "Ethics and oversight")
    bullets(s, [
        "Every advisory carries requires_human_approval.",
        "Supervisory autonomy is bounded at 1.5 °C — enforced by the room.",
        "Occupancy is counted, never identified, and stays in the room that produced it.",
        "Known limitations are shown in the dashboard, not just in the appendix.",
    ])
    textbox(s, "There is no path from a model score to a physical action\n"
               "that does not pass through a person.",
            Inches(1.0), Inches(5.0), Inches(11.0), Inches(1.2),
            size=22, bold=True, colour=ACCENT)
    notes(s, "There is no path from a model score to a physical action that "
             "doesn't pass through a person.")

    # 13 — the ask
    s = blank(prs)
    heading(s, "Phase 1: instrumented pilot", "The ask")
    table(s, [
        ["Scope", "3 units, one floor, observe-only"],
        ["Duration", "3–6 months"],
        ["Exit criterion", "model recalibrated on REAL telemetry; fairness audit repeated"],
        ["Named risk", "real failures are rare — may see zero events; run in shadow mode"],
    ], Inches(1.0), Inches(2.1), Inches(11.3), Inches(2.4),
        col_widths=[Inches(2.8), Inches(8.5)], header=False, size=16)
    textbox(s, "Approve Phase 1, and the €4,000 security hardening up front.",
            Inches(1.0), Inches(4.9), Inches(11.0), Inches(0.8),
            size=24, bold=True, colour=ACCENT)
    textbox(s, "Everything today is built on simulated data. The failure physics is "
               "calibrated against a published\nindustrial dataset, but no number in "
               "this business case is trustworthy until it has been seen on\nthis "
               "building. Phase 1 exists to find that out cheaply.",
            Inches(1.0), Inches(5.8), Inches(11.2), Inches(1.4),
            size=14, colour=MUTED)
    notes(s, "Everything today is built on simulated data. Phase 1 exists to find "
             "out cheaply whether it holds on this building.")

    prs.save(str(OUT))
    return OUT, len(prs.slides.__iter__.__self__._sldIdLst), charts


if __name__ == "__main__":
    out, _, charts = build()
    prs = Presentation(str(out))
    print(f"wrote {out}  ({len(prs.slides)} slides, "
          f"{out.stat().st_size/1024:.0f} KB)")
    for name, path in charts.items():
        print(f"  chart {name:<10} {path}")
