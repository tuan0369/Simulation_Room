"""Shared chart and slide helpers for the Project-2 decks.

Both `report/build_report.py` (the project report) and
`report/pitch/build_pitch.py` (the executive pitch deliverable) draw from here,
so a retrained model updates both and they cannot disagree about the numbers.

Charts are generated from `ml/models/*.csv` rather than drawn by hand.
"""
from __future__ import annotations

import json
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt          # noqa: E402
import pandas as pd                      # noqa: E402
from pptx.dml.color import RGBColor      # noqa: E402
from pptx.util import Inches, Pt         # noqa: E402

ROOT = Path(".")
MODELS = ROOT / "ml" / "models"
ASSETS = ROOT / "report" / "assets"
DEMO = ROOT / "report" / "demo"

INK = RGBColor(0x0F, 0x17, 0x2A)
MUTED = RGBColor(0x47, 0x55, 0x69)
ACCENT = RGBColor(0x25, 0x63, 0xEB)
WARN = RGBColor(0xB4, 0x53, 0x09)
GOOD = RGBColor(0x15, 0x80, 0x3D)

W, H = Inches(13.333), Inches(7.5)

plt.rcParams.update({
    "figure.dpi": 200, "savefig.bbox": "tight", "font.size": 11,
    "axes.spines.top": False, "axes.spines.right": False,
    "axes.grid": True, "grid.alpha": 0.25, "axes.axisbelow": True,
})

# ROI figures — single source, mirrored in docs/roi_roadmap.md §3.
SCALE = [(6, 5.7), (12, 2.2), (24, 1.2), (50, 0.8)]
BENEFITS = [("Fan energy\n(measured)", 155),
            ("Avoided outages\n(assumed rate)", 3974),
            ("Condition-based\nservicing", 1800)]

FAULT_LABELS = {"hdf": "Heat dissipation", "osf": "Overstrain", "pwf": "Power",
                "airflow": "Airflow", "bearing": "Bearing"}


def spec() -> dict:
    return json.loads((MODELS / "feature_spec.json").read_text())


# ── Charts ──────────────────────────────────────────────────────────────────

def chart_model_comparison() -> Path:
    df = pd.read_csv(MODELS / "model_comparison.csv")
    labels = {"always 'no failure'": "Always\n\"no failure\"",
              "always predict no-failure": "Always\n\"no failure\"",
              "threshold motor_temp>80": "Single\nthreshold",
              "logistic_regression": "Logistic\nregression",
              "random_forest": "Random\nforest",
              "gradient_boosting": "Gradient boosting\n(shipped)"}
    df = df[df.model.isin(labels)]
    names = [labels[m] for m in df.model]
    colours = ["#2563eb" if "shipped" in n else "#94a3b8" for n in names]
    fig, ax = plt.subplots(figsize=(7.8, 3.4))
    bars = ax.bar(names, df.pr_auc, color=colours, width=0.62)
    for b, v in zip(bars, df.pr_auc):
        ax.text(b.get_x() + b.get_width()/2, v + 0.02, f"{v:.2f}",
                ha="center", fontsize=10, weight="bold")
    ax.tick_params(axis="x", labelsize=9.5)
    ax.set_ylabel("PR-AUC"); ax.set_ylim(0, 1.08)
    ax.set_title("Ten times better than the obvious rule", loc="left",
                 weight="bold")
    return _save(fig, "pitch_model_comparison")


def chart_fault_modes() -> Path:
    df = pd.read_csv(MODELS / "recall_by_fault_mode.csv")
    df["label"] = df.fault.map(FAULT_LABELS).fillna(df.fault)
    x = range(len(df))
    fig, ax = plt.subplots(figsize=(7.4, 3.2))
    ax.bar([i - 0.2 for i in x], df.recall_model, 0.4,
           label="Model alone", color="#94a3b8")
    ax.bar([i + 0.2 for i in x], df.recall_hybrid, 0.4,
           label="With thermal guard", color="#2563eb")
    ax.set_xticks(list(x)); ax.set_xticklabels(df.label, fontsize=9)
    ax.set_ylabel("Recall"); ax.set_ylim(0, 1.15); ax.legend(frameon=False)
    ax.set_title("All five failure modes detected", loc="left", weight="bold")
    return _save(fig, "pitch_fault_modes")


def chart_fairness() -> Path:
    df = pd.read_csv(MODELS / "fairness_audit.csv").dropna(subset=["recall"])
    df = df.sort_values("recall")
    names = [t.split("/")[-1].replace("-", " ").title() for t in df.twin_id]
    fig, ax = plt.subplots(figsize=(7.4, 3.2))
    ax.barh([i + 0.2 for i in range(len(df))], df.recall, 0.4,
            label="Recall", color="#16a34a")
    ax.barh([i - 0.2 for i in range(len(df))], df.precision, 0.4,
            label="Precision", color="#f59e0b")
    ax.set_yticks(range(len(df))); ax.set_yticklabels(names)
    ax.set_xlim(0, 1.05); ax.legend(frameon=False, loc="lower right")
    ax.set_title("Detection is even; precision is not", loc="left",
                 weight="bold")
    return _save(fig, "pitch_fairness")


def chart_payback() -> Path:
    units = [u for u, _ in SCALE]
    years = [y for _, y in SCALE]
    fig, ax = plt.subplots(figsize=(7.4, 3.2))
    bars = ax.bar([str(u) for u in units], years,
                  color=["#dc2626" if y > 3 else "#16a34a" for y in years])
    for b, y in zip(bars, years):
        ax.text(b.get_x() + b.get_width()/2, y + 0.12, f"{y} yr",
                ha="center", fontsize=10, weight="bold")
    ax.axhline(3, color="#475569", ls="--", lw=1)
    ax.text(3.35, 3.12, "3-year hurdle", fontsize=9, color="#475569")
    ax.set_xlabel("HVAC units in scope"); ax.set_ylabel("Payback (years)")
    ax.set_ylim(0, 7)
    ax.set_title("Scale, not tuning, is what makes this pay", loc="left",
                 weight="bold")
    return _save(fig, "pitch_payback")


def chart_benefits() -> Path:
    labels = [n for n, _ in BENEFITS]
    values = [v for _, v in BENEFITS]
    fig, ax = plt.subplots(figsize=(7.4, 3.0))
    bars = ax.bar(labels, values, color=["#94a3b8", "#2563eb", "#60a5fa"])
    for b, v in zip(bars, values):
        ax.text(b.get_x() + b.get_width()/2, v + 90, f"€{v:,}",
                ha="center", fontsize=10, weight="bold")
    ax.set_ylabel("€ per year"); ax.set_ylim(0, 4800)
    ax.set_title("The energy saving is real, and small", loc="left",
                 weight="bold")
    return _save(fig, "pitch_benefits")


def chart_room_faults() -> Path:
    """Which fault each unit is prone to — the ecosystem's variety."""
    df = pd.read_csv(MODELS / "fairness_audit.csv").dropna(subset=["recall"])
    names = [t.split("/")[-1].replace("-", " ").title() for t in df.twin_id]
    fig, ax = plt.subplots(figsize=(7.4, 3.0))
    bars = ax.bar(names, df.positive_rate_pct, color="#3b82f6")
    for b, v in zip(bars, df.positive_rate_pct):
        ax.text(b.get_x() + b.get_width()/2, v + 0.06, f"{v:.1f}%",
                ha="center", fontsize=10, weight="bold")
    ax.set_ylabel("Rows within 4 h of a failure")
    ax.tick_params(axis="x", labelsize=9)
    ax.set_title("Each unit has its own wear character", loc="left",
                 weight="bold")
    return _save(fig, "report_room_faults")


def _save(fig, name: str) -> Path:
    ASSETS.mkdir(parents=True, exist_ok=True)
    path = ASSETS / f"{name}.png"
    fig.savefig(path)
    plt.close(fig)
    return path


def all_charts() -> dict:
    return {
        "models": chart_model_comparison(),
        "modes": chart_fault_modes(),
        "fairness": chart_fairness(),
        "payback": chart_payback(),
        "benefits": chart_benefits(),
        "rooms": chart_room_faults(),
    }


# ── Slide helpers ───────────────────────────────────────────────────────────

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def textbox(slide, text, left, top, width, height, size=18, bold=False,
            colour=INK, align=None, spacing=1.15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.line_spacing = spacing
        if align is not None:
            p.alignment = align
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = colour
            run.font.name = "Calibri"
    return box


def heading(slide, title, kicker=None):
    if kicker:
        textbox(slide, kicker.upper(), Inches(0.8), Inches(0.45), Inches(11),
                Inches(0.4), size=12, bold=True, colour=ACCENT)
    textbox(slide, title, Inches(0.8), Inches(0.8), Inches(11.8), Inches(1.0),
            size=30, bold=True)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def bullets(slide, items, top=Inches(2.0), size=18, colour=INK, left=Inches(1.0),
            width=Inches(11.4), step=Inches(0.62)):
    y = top
    for item in items:
        textbox(slide, "•  " + item, left, y, width, Inches(0.6),
                size=size, colour=colour)
        y += step
    return y


def picture(slide, path, left, top, width):
    if Path(path).exists():
        slide.shapes.add_picture(str(path), left, top, width=width)


def table(slide, rows, left, top, width, height, col_widths=None,
          header=True, size=15):
    shape = slide.shapes.add_table(len(rows), len(rows[0]), left, top, width,
                                   height)
    tbl = shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = w
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(size)
                    run.font.bold = bool(header and r == 0)
                    run.font.name = "Calibri"
    return tbl
