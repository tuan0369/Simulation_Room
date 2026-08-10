"""Build the Project 2 report deck.

    docker compose run --rm sim python report/build_report.py

Structured like Project 1's report — objective, then a walk through the brief,
ending in a deliverables checklist — but against the Project 2 requirements:

    Key tasks       Predictive Intelligence · Ecosystem Integration ·
                    Governance & Ethics · Strategic Roadmap
    Deliverables    Predictive Model Output · Integrated Ecosystem Diagram ·
                    Executive Pitch

Charts come from `ml/models/*.csv` via `deckkit`, shared with the executive
pitch, so a retrained model updates both decks and they cannot disagree.
"""
from __future__ import annotations

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

from pptx import Presentation                                   # noqa: E402
from pptx.util import Inches                                    # noqa: E402

from deckkit import (ACCENT, DEMO, GOOD, INK, MUTED, W, H, WARN,  # noqa: E402
                     all_charts, blank, bullets, heading, notes, picture,
                     spec, table, textbox)

OUT = Path("report/Digital_Twin_Project2_Report.pptx")


def build():
    charts = all_charts()
    m = spec()["metrics"]
    thr = spec()["decision_threshold"]

    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    # 1 — title
    s = blank(prs)
    textbox(s, "Smart Facility Digital Twin", Inches(0.9), Inches(2.1),
            Inches(11.5), Inches(1.2), size=44, bold=True)
    textbox(s, "Project 2 — Intelligent Ecosystem & Strategic Optimization",
            Inches(0.9), Inches(3.3), Inches(11.5), Inches(0.7), size=22,
            colour=MUTED)
    textbox(s, "2 floors · 6 rooms · 11 interacting twins · "
               "5 predicted failure modes",
            Inches(0.9), Inches(4.1), Inches(11.5), Inches(0.6), size=16,
            colour=ACCENT)
    textbox(s, "Extends the Project 1 single-room twin into a coordinated, "
               "self-monitoring facility.",
            Inches(0.9), Inches(5.0), Inches(11.5), Inches(0.6), size=14,
            colour=MUTED)
    notes(s, "Project 2 report. Everything shown runs; nothing is mocked.")

    # 2 — objective and scope
    s = blank(prs)
    heading(s, "Objective & scope", "Brief")
    bullets(s, [
        "Extend one room twin into an intelligent, autonomous ecosystem.",
        "Integrate a predictive model for equipment failure.",
        "Coordinate multiple interacting twins across a smart facility.",
        "Address governance, ethics, ROI and a deployment roadmap.",
    ], top=Inches(1.9))
    table(s, [
        ["Key task", "Delivered by"],
        ["Predictive Intelligence", "5-class fault model + physics guard, live per room"],
        ["Ecosystem Integration", "11 twins, federated coordination, MQTT data flows"],
        ["Governance & Ethics", "Threat model, fairness audit, bounded autonomy"],
        ["Strategic Roadmap", "Measured ROI, 4-phase pilot-to-campus plan"],
    ], Inches(1.0), Inches(4.5), Inches(11.3), Inches(2.2),
        col_widths=[Inches(3.4), Inches(7.9)])
    notes(s, "Four key tasks from the brief, and where each is answered.")

    # 3 — from one room to a facility
    s = blank(prs)
    heading(s, "From one room to a facility", "What changed")
    table(s, [
        ["", "Project 1", "Project 2"],
        ["Scope", "1 room", "2 floors, 6 rooms, 11 twins"],
        ["Coordination", "None", "Federated — rooms decide, supervisors advise"],
        ["Occupancy", "Per-room random walk", "Headcount-conserving people flow"],
        ["Equipment", "Never degrades", "5 failure modes, AI4I-calibrated"],
        ["Intelligence", "PID control", "PID + fault-type model + thermal guard"],
        ["Remediation", "AC on/off", "5 remedies, manual or model-dispatched"],
    ], Inches(1.0), Inches(2.0), Inches(11.3), Inches(3.6),
        col_widths=[Inches(2.4), Inches(3.9), Inches(5.0)])
    textbox(s, "The Project-1 room survives as f1/lab-a with its calibrated "
               "constants intact — a test locks them as the regression guard.",
            Inches(1.0), Inches(5.9), Inches(11.0), Inches(0.8), size=15,
            colour=MUTED)
    notes(s, "Project 1's room is still here, unchanged, as the regression "
             "guard for the whole refactor.")

    # 4 — ecosystem integration
    s = blank(prs)
    heading(s, "Eleven twins, each with one job", "Ecosystem integration")
    table(s, [
        ["Twin", "Count", "Owns", "Loop"],
        ["Room twin", "6", "Physics, PID, equipment health", "1 s"],
        ["Floor twin", "2", "Aggregate; recommend setpoint nudges", "10 s"],
        ["Building twin", "1", "kW allocation; maintenance work orders", "30 s"],
        ["Occupancy twin", "1", "People flow between rooms", "1 s"],
        ["Risk scorer", "1", "Fault probability per room", "30 s"],
    ], Inches(1.0), Inches(2.0), Inches(11.3), Inches(3.0),
        col_widths=[Inches(2.6), Inches(1.1), Inches(5.6), Inches(2.0)])
    textbox(s, "Occupancy → energy → equipment wear → maintenance is a single "
               "causal chain.\nA class leaving one room raises load in another, "
               "which accrues wear, which the model sees hours later.",
            Inches(1.0), Inches(5.3), Inches(11.2), Inches(1.2), size=16,
            colour=ACCENT)
    notes(s, "The brief asks for an energy twin interacting with an occupancy "
             "twin. That path is real: headcount is conserved, so people "
             "genuinely move between rooms and the thermal load follows them.")

    # 5 — coordination strategy (required deliverable)
    s = blank(prs)
    heading(s, "Centralised vs federated", "Deliverable — coordination strategy")
    table(s, [
        ["", "Centralised", "Federated (chosen)"],
        ["A crash stops…", "every room cooling", "coordination only"],
        ["Local disturbance", "waits for the global loop", "handled in the room's own 1 s loop"],
        ["Occupancy data", "leaves every room", "stays room-local; aggregates only"],
        ["Tuning", "one loop for a 20 m² and a 70 m² room", "per-room thermal mass"],
        ["A supervisor bug", "actuates every room", "bounded advice, ≤ 1.5 °C"],
    ], Inches(1.0), Inches(2.0), Inches(11.3), Inches(3.2),
        col_widths=[Inches(2.6), Inches(4.2), Inches(4.5)])
    textbox(s, "The 1.5 °C limit is enforced by the room, not by the supervisor "
               "sending the advice — so a spoofed or buggy advisory cannot make "
               "a room unsafe. Eight named tests back these claims.",
            Inches(1.0), Inches(5.5), Inches(11.2), Inches(1.0), size=16,
            colour=ACCENT)
    notes(s, "Federated with hierarchical supervision. Authority flows downward "
             "as advice and stops at a bound the recipient owns.")

    # 6 — the failure modes
    s = blank(prs)
    heading(s, "Five failure modes, one per unit character",
            "Predictive intelligence")
    picture(s, charts["rooms"], Inches(0.8), Inches(1.9), Inches(7.2))
    table(s, [
        ["Unit", "Fails by"],
        ["Wet lab A", "Airflow"],
        ["Dry lab B", "Bearing"],
        ["Server room", "Overstrain"],
        ["Teaching lab C", "Heat dissipation"],
        ["Meeting room", "Power"],
    ], Inches(8.4), Inches(2.1), Inches(4.3), Inches(2.6),
        col_widths=[Inches(2.3), Inches(2.0)], size=14)
    textbox(s, "Thresholds inherit from UCI AI4I 2020, whose published rules we "
               "verified reproduce exactly (115/115 rows).",
            Inches(0.8), Inches(5.6), Inches(11.4), Inches(0.8), size=15,
            colour=MUTED)
    notes(s, "Real buildings do not contain six identical units. Each has its "
             "own wear character, which is what makes fault classification a "
             "real problem rather than one class with noise.")

    # 7 — model results
    s = blank(prs)
    heading(s, "Model results", "Deliverable — predictive model output")
    picture(s, charts["models"], Inches(0.8), Inches(1.9), Inches(7.4))
    textbox(s, f"PR-AUC  {m['pr_auc']:.3f}\nRecall  {m['recall']:.3f}\n"
               f"Precision  {m['precision']:.3f}",
            Inches(8.6), Inches(2.3), Inches(4.0), Inches(2.0), size=22,
            bold=True)
    textbox(s, f"Decision threshold {thr:.4f}, chosen from a\ncost curve — not "
               f"left at 0.5.\n\nAccuracy would be 97 % if we always\npredicted "
               f"\"no failure\". That is why\nit is not the headline.",
            Inches(8.6), Inches(4.3), Inches(4.2), Inches(2.0), size=14,
            colour=MUTED)
    notes(s, "Trained and evaluated with a temporal split and a held-out room. "
             "A random shuffle would inflate this to a meaningless 0.99.")

    # 8 — per mode, and the correction
    s = blank(prs)
    heading(s, "Every mode detected — after a correction", "Predictive intelligence")
    picture(s, charts["modes"], Inches(0.8), Inches(1.9), Inches(7.4))
    textbox(s, "Heat dissipation:\n0.00 → 0.93 recall",
            Inches(8.6), Inches(2.2), Inches(4.2), Inches(1.1), size=20,
            bold=True, colour=WARN)
    textbox(s, "The first model was blind to that mode.\n"
               "It had seen about two examples of it.\n\n"
               "The cause was starvation, not tuning:\n"
               "six identical units meant one mode\n"
               "was 90 % of the data. Fixing the DATA\n"
               "fixed the model.\n\n"
               "A physics threshold still backs it up.",
            Inches(8.6), Inches(3.4), Inches(4.3), Inches(2.8), size=13,
            colour=MUTED)
    notes(s, "A model that cannot see a failure mode is usually being starved "
             "of it, not badly tuned. Worth stating because it generalises.")

    # 9 — fairness
    s = blank(prs)
    heading(s, "Fairness audit", "Governance & ethics")
    picture(s, charts["fairness"], Inches(0.8), Inches(1.9), Inches(7.4))
    textbox(s, "Recall 0.89–1.00\nin every room",
            Inches(8.6), Inches(2.2), Inches(4.2), Inches(1.1), size=20,
            bold=True, colour=GOOD)
    textbox(s, "Room identity is excluded from the\nfeatures by construction, so "
               "the model\ncannot learn a per-room prior.\n\n"
               "Teaching Lab C stays weakest:\nprecision 0.45, and it misses "
               "about\none failure in nine.\n\n"
               "The audit is repeated after every retrain.",
            Inches(8.6), Inches(3.4), Inches(4.3), Inches(2.8), size=13,
            colour=MUTED)
    notes(s, "An earlier version had a 100 % false-negative rate on the wet lab "
             "— the room where an outage costs most. We publish both states.")

    # 10 — remediation
    s = blank(prs)
    heading(s, "Closing the loop: one remedy per mode", "Ecosystem integration")
    table(s, [
        ["Fault", "Remedy", "Reduces cooling?"],
        ["Airflow", "Replace filter", "No"],
        ["Bearing", "Service motor", "No"],
        ["Overstrain", "Service motor + occupant notice", "No"],
        ["Power", "Electrical service", "No"],
        ["Heat dissipation", "Thermal derate — cap fan duty", "Yes, bounded"],
    ], Inches(1.0), Inches(2.0), Inches(11.3), Inches(3.0),
        col_widths=[Inches(2.8), Inches(5.5), Inches(3.0)])
    textbox(s, "The derate is capped at 50 % duty, never zero, and releases "
               "itself once the winding cools —\nwhat a thermal overload relay "
               "does. Easing a motor off is not the same act as switching "
               "cooling off.",
            Inches(1.0), Inches(5.3), Inches(11.2), Inches(1.2), size=15,
            colour=MUTED)
    notes(s, "The system used to predict five faults and fix two. Three "
             "predictions had no action anyone could take.")

    # 11 — security
    s = blank(prs)
    heading(s, "Security — the honest posture", "Governance & ethics")
    bullets(s, [
        "The pilot broker is anonymous and unencrypted. Fine on one machine.",
        "Unacceptable in a building that dispatches technicians to equipment.",
    ], top=Inches(1.9))
    table(s, [
        ["Hardening", "Why"],
        ["mTLS + per-topic ACLs", "a room twin can only speak for itself"],
        ["Signed commands", "no replay, no spoofed shutdown"],
        ["Append-only audit log", "every action attributable to a person"],
        ["Network segmentation", "not routable from the office LAN"],
    ], Inches(1.0), Inches(3.3), Inches(11.3), Inches(2.4),
        col_widths=[Inches(4.3), Inches(7.0)])
    textbox(s, "Costed at €4,000 — a line item, not a backlog ticket.",
            Inches(1.0), Inches(6.0), Inches(9.0), Inches(0.6), size=17,
            bold=True, colour=ACCENT)
    notes(s, "Stating the pilot is insecure is more useful than claiming "
             "otherwise. The threat model and checklist are in governance.md.")

    # 12 — autonomy
    s = blank(prs)
    heading(s, "Bounded autonomy, in one operator control", "Governance & ethics")
    table(s, [
        ["Level", "Thermostat runs the AC", "Model dispatches maintenance"],
        ["Manual", "No", "No"],
        ["Auto climate  (default)", "Yes", "No — approval required"],
        ["Full auto", "Yes", "Yes — preventive actions only"],
    ], Inches(1.0), Inches(2.0), Inches(11.3), Inches(2.2),
        col_widths=[Inches(3.5), Inches(3.9), Inches(3.9)])
    bullets(s, [
        "Full auto is never the default — autonomy is opted into.",
        "Preventive actions only; at most one automatic service per unit per 24 h.",
        "Every dispatch is still published as an advisory, so the audit trail is identical.",
        "Shutting a unit down is forbidden at every level.",
    ], top=Inches(4.5), size=16)
    notes(s, "There is no path from a model score to a physical action that "
             "does not pass through a person, or through these bounds.")

    # 13 — ROI
    s = blank(prs)
    heading(s, "Where the value is — and is not", "Strategic roadmap")
    picture(s, charts["benefits"], Inches(0.8), Inches(1.9), Inches(7.4))
    textbox(s, "≈ €5,929 / year", Inches(8.7), Inches(2.4), Inches(4.0),
            Inches(0.8), size=26, bold=True, colour=GOOD)
    textbox(s, "Fan energy is measured: 620 kWh, €155.\n"
               "Small. Anyone pitching predictive\n"
               "maintenance on fan energy alone is\n"
               "selling you something.\n\n"
               "The value is in outages you don't have.",
            Inches(8.7), Inches(3.4), Inches(4.2), Inches(2.2), size=14,
            colour=MUTED)
    notes(s, "Every figure is labelled measured or assumed. The cooling-capacity "
             "loss is probably larger and we do not claim it, because we did "
             "not measure it.")

    # 14 — the uncomfortable number
    s = blank(prs)
    heading(s, "At six units it does not pay for itself", "Strategic roadmap")
    picture(s, charts["payback"], Inches(0.8), Inches(1.9), Inches(7.4))
    table(s, [
        ["One-off", "€17,000"],
        ["Recurring", "€2,900 / yr"],
        ["Payback", "5.7 years"],
        ["5-yr NPV @ 8 %", "−€5,021"],
    ], Inches(8.7), Inches(2.2), Inches(4.0), Inches(1.9),
        col_widths=[Inches(2.2), Inches(1.8)], header=False, size=15)
    textbox(s, "Integration cost is roughly fixed;\nbenefits scale with unit "
               "count.\n\nRecommendation: pilot to prove the\nmodel on real "
               "telemetry, then scale\nto campus — not building by building.",
            Inches(8.7), Inches(4.3), Inches(4.2), Inches(2.2), size=14,
            colour=WARN)
    notes(s, "Reporting a negative NPV is the point. The conclusion is robust "
             "across the sensitivity table.")

    # 15 — roadmap
    s = blank(prs)
    heading(s, "Pilot to campus", "Strategic roadmap")
    table(s, [
        ["Phase", "Scope", "Exit criterion", "Named risk"],
        ["0 — Pilot\n(complete)", "6 simulated twins", "Suite green; loop verified end to end", "All results rest on simulated data"],
        ["1 — Instrumented\n3–6 months", "3 real units, observe-only", "Model recalibrated on REAL telemetry", "Real failures are rare; may see zero"],
        ["2 — Full building\n6–12 months", "All units, live work orders", "Measured drop in unplanned outages", "Alert fatigue at 0.73 precision"],
        ["3 — Campus\n12–24 months", "≥ 24 units", "Positive measured NPV", "Models transfer poorly (PR-AUC 0.20)"],
    ], Inches(0.8), Inches(1.9), Inches(11.8), Inches(3.6),
        col_widths=[Inches(2.3), Inches(2.7), Inches(3.5), Inches(3.3)], size=12)
    textbox(s, "Phase 4 — bounded autonomy — only after Phase 3, and never "
               "automatic shutdown.",
            Inches(0.8), Inches(5.9), Inches(11.4), Inches(0.6), size=16,
            colour=ACCENT)
    notes(s, "Each phase has an entry criterion, an exit criterion and a named "
             "risk with its mitigation.")

    # 16 — how to run
    s = blank(prs)
    heading(s, "Working prototype — how to run", "Reproducibility")
    textbox(s,
            "docker compose up -d\n"
            "    dashboard   http://localhost:8501\n"
            "    3D view     http://localhost:8000/room3d/building3d.html\n\n"
            "docker compose run --rm sim pytest -v\n"
            "docker compose run --rm sim python simulator/dataset_generator.py\n"
            "docker compose run --rm sim python ml/train.py",
            Inches(1.0), Inches(2.0), Inches(11.0), Inches(2.8), size=16)
    bullets(s, [
        "Everything runs in Docker — no host Python needed.",
        "The dataset and both models regenerate from source, deterministically.",
        "Notebooks in ml/notebooks/ carry executed output.",
    ], top=Inches(5.0), size=16, colour=MUTED)
    demo = sorted(DEMO.glob("*.png"))
    notes(s, "One command brings up the broker, simulator, dashboard and 3D "
             "view. The full suite runs in the same image.")

    # 17 — deliverables checklist
    s = blank(prs)
    heading(s, "Deliverables", "Conclusion")
    table(s, [
        ["Required deliverable", "Where"],
        ["Predictive model output — notebooks / results",
         "ml/notebooks/*.ipynb · ml/models/model_card.md"],
        ["Integrated ecosystem diagram — coordination strategy",
         "docs/ecosystem.md (data flow, sequence, centralised vs federated)"],
        ["Executive pitch — ROI, security, ethical safety",
         "report/pitch/Digital_Twin_Project2_Executive_Pitch.pptx"],
        ["Predictive intelligence (key task)",
         "5-class model + thermal guard, live per room"],
        ["Ecosystem integration (key task)",
         "11 twins, federated coordination, MQTT topic map"],
        ["Governance & ethics (key task)", "docs/governance.md"],
        ["Strategic roadmap (key task)", "docs/roi_roadmap.md"],
    ], Inches(0.8), Inches(1.9), Inches(11.8), Inches(3.9),
        col_widths=[Inches(5.2), Inches(6.6)], size=13)
    textbox(s, "All results are reproducible from a clean clone. "
               "The models are trained on simulated telemetry — "
               "recalibration on real data is Phase 1.",
            Inches(0.8), Inches(6.1), Inches(11.6), Inches(0.8), size=15,
            colour=MUTED)
    notes(s, "Closing slide: every required deliverable and where to find it, "
             "with the standing caveat about simulated data.")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    return OUT, charts


if __name__ == "__main__":
    out, charts = build()
    prs = Presentation(str(out))
    print(f"wrote {out}  ({len(prs.slides)} slides, "
          f"{out.stat().st_size/1024:.0f} KB)")
