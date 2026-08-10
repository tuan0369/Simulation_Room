"""Checks on the Project 2 report deck.

The report walks the brief's four key tasks and three deliverables, so these
verify that structure is actually present — and that the report and the
executive pitch, which share their chart code, agree on the numbers.
"""
from pathlib import Path

import pytest

pptx = pytest.importorskip("pptx")
from pptx import Presentation  # noqa: E402

REPORT = Path("report/Digital_Twin_Project2_Report.pptx")
PITCH = Path("report/pitch/Digital_Twin_Project2_Executive_Pitch.pptx")

pytestmark = pytest.mark.skipif(
    not REPORT.exists(), reason="report not built; run report/build_report.py")


def _text(path):
    deck = Presentation(str(path))
    out = []
    for slide in deck.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                out.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    out.extend(c.text for c in row.cells)
        out.append(slide.notes_slide.notes_text_frame.text)
    return "\n".join(out)


@pytest.fixture(scope="module")
def deck():
    return Presentation(str(REPORT))


@pytest.fixture(scope="module")
def text():
    return _text(REPORT)


def test_report_is_a_sensible_length(deck):
    assert 12 <= len(deck.slides) <= 22


def test_widescreen_like_project_one(deck):
    assert round(deck.slide_width / deck.slide_height, 2) == pytest.approx(
        1.78, abs=0.02)


def test_every_slide_has_speaker_notes(deck):
    for i, slide in enumerate(deck.slides, 1):
        assert slide.notes_slide.notes_text_frame.text.strip(), (
            f"slide {i} has no speaker notes")


def test_charts_are_embedded(deck):
    pictures = sum(1 for s in deck.slides for sh in s.shapes
                   if sh.shape_type == 13)
    assert pictures >= 5


# ── The brief's structure must be present ───────────────────────────────────

@pytest.mark.parametrize("task", [
    "Predictive Intelligence", "Ecosystem Integration",
    "Governance & Ethics", "Strategic Roadmap",
])
def test_every_key_task_from_the_brief_appears(text, task):
    assert task.lower() in text.lower(), f"key task missing: {task}"


@pytest.mark.parametrize("deliverable", [
    "predictive model output", "ecosystem diagram", "executive pitch",
])
def test_every_required_deliverable_is_named(text, deliverable):
    assert deliverable.lower() in text.lower(), (
        f"deliverable missing: {deliverable}")


def test_report_follows_the_project_one_shape(text):
    """Objective, a walk through the brief, then a deliverables checklist."""
    assert "Objective & scope" in text
    assert "how to run" in text.lower()
    assert "Deliverables" in text


def test_deliverables_checklist_points_at_real_files(text):
    for path in ("ml/notebooks", "docs/ecosystem.md", "docs/governance.md",
                 "docs/roi_roadmap.md"):
        assert path in text, f"checklist does not cite {path}"
        assert Path(path.split()[0]).exists(), f"{path} does not exist"


def test_centralised_vs_federated_is_covered(text):
    """The brief asks specifically for the coordination strategy."""
    assert "federated" in text.lower()
    assert "centralised" in text.lower() or "centralized" in text.lower()


# ── Honesty, same standard as the pitch ─────────────────────────────────────

def test_report_admits_the_data_is_simulated(text):
    assert "simulated" in text.lower()


def test_report_states_the_negative_npv(text):
    assert "does not pay for itself" in text.lower()
    assert "−€5,021" in text or "-€5,021" in text


def test_report_states_the_security_posture(text):
    assert "anonymous and unencrypted" in text.lower()


def test_report_keeps_the_fairness_finding(text):
    assert "fairness" in text.lower()
    assert "precision 0.45" in text or "0.45" in text


def test_report_records_the_blind_spot_correction(text):
    assert "0.00" in text and "0.93" in text


# ── Consistency between the two decks ───────────────────────────────────────

@pytest.mark.skipif(not PITCH.exists(), reason="pitch not built")
def test_report_and_pitch_agree_on_the_headline_metrics():
    """They share deckkit precisely so they cannot drift apart."""
    import json
    m = json.loads(Path("ml/models/feature_spec.json").read_text())["metrics"]
    for body in (_text(REPORT), _text(PITCH)):
        assert f"{m['pr_auc']:.2f}" in body or f"{m['pr_auc']:.3f}" in body
        assert f"{m['recall']:.2f}" in body or f"{m['recall']:.3f}" in body


@pytest.mark.skipif(not PITCH.exists(), reason="pitch not built")
def test_both_decks_share_one_chart_source():
    pitch_src = Path("report/pitch/build_pitch.py").read_text(encoding="utf-8")
    report_src = Path("report/build_report.py").read_text(encoding="utf-8")
    for src in (pitch_src, report_src):
        assert "from deckkit import" in src
        assert "def chart_" not in src, "a deck redefines its own charts"
