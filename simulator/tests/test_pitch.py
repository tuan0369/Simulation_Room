"""Checks on the executive pitch deck.

The deck is generated from the project's own artifacts, so these verify it was
actually built from them — and that the honesty slides survived. A pitch that
quietly drops the limitations is the failure mode worth testing for.
"""
from pathlib import Path

import pytest

pptx = pytest.importorskip("pptx")
from pptx import Presentation  # noqa: E402

DECK = Path("report/pitch/Digital_Twin_Project2_Executive_Pitch.pptx")
OUTLINE = Path("report/pitch/pitch_outline.md")

pytestmark = pytest.mark.skipif(not DECK.exists(),
                                reason="deck not built yet; run build_pitch.py")


@pytest.fixture(scope="module")
def deck():
    return Presentation(str(DECK))


@pytest.fixture(scope="module")
def text(deck):
    out = []
    for slide in deck.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                out.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    out.extend(c.text for c in row.cells)
        out.append(slide.notes_slide.notes_text_frame.text)
    return "\n".join(out)


def test_deck_has_a_sensible_length(deck):
    assert 10 <= len(deck.slides) <= 16


def test_widescreen(deck):
    assert round(deck.slide_width / deck.slide_height, 2) == pytest.approx(1.78, abs=0.02)


def test_every_slide_has_speaker_notes(deck):
    for i, slide in enumerate(deck.slides, 1):
        assert slide.notes_slide.notes_text_frame.text.strip(), (
            f"slide {i} has no speaker notes")


def test_charts_are_embedded(deck):
    pictures = sum(1 for s in deck.slides for sh in s.shapes if sh.shape_type == 13)
    assert pictures >= 5, "generated charts are missing from the deck"


def test_charts_were_generated_from_artifacts():
    for name in ("model_comparison", "fault_modes", "fairness", "payback",
                 "benefits"):
        assert Path(f"report/assets/pitch_{name}.png").exists()


# ── The honesty slides must survive ─────────────────────────────────────────

def test_deck_states_the_model_blind_spot(text):
    assert "cannot do" in text or "Blind to" in text
    assert "heat-dissipation" in text.lower() or "heat dissipation" in text.lower()


def test_deck_states_the_fairness_finding(text):
    assert "Wet Lab A" in text
    assert "100 %" in text or "100%" in text


def test_deck_states_the_negative_npv(text):
    """The business case is negative at six units. A pitch that hides that is
    the exact failure this test exists to prevent."""
    assert "does not pay for itself" in text
    assert "−€5,420" in text or "-€5,420" in text


def test_deck_admits_the_data_is_simulated(text):
    assert "simulated data" in text


def test_deck_states_the_security_posture(text):
    assert "not secure" in text


def test_deck_states_human_approval(text):
    assert "never switches anything off" in text or "requires_human_approval" in text


def test_deck_does_not_lead_with_accuracy(text):
    """At a 1.7% positive rate, accuracy is a misleading headline and the deck
    should say so rather than quote it."""
    if "98 %" in text or "98%" in text:
        assert "why we don" in text.lower() or "that is why" in text.lower()


# ── Consistency with the sources ────────────────────────────────────────────

def test_metrics_match_the_trained_model(text):
    import json
    spec = json.loads(Path("ml/models/feature_spec.json").read_text())
    m = spec["metrics"]
    assert f"{m['pr_auc']:.2f}" in text
    assert f"{m['recall']:.2f}" in text


def test_roi_figures_match_the_roadmap(text):
    roadmap = Path("docs/roi_roadmap.md").read_text(encoding="utf-8")
    for figure in ("17,000", "2,900", "5.9"):
        assert figure in text, f"{figure} missing from the deck"
        assert figure in roadmap, f"{figure} missing from the roadmap"


def test_outline_exists_and_traces_its_numbers():
    body = OUTLINE.read_text(encoding="utf-8")
    assert "model_card.md" in body and "roi_roadmap.md" in body
    assert "Numbers appearing in the deck" in body
